from fastapi import FastAPI, UploadFile, File, HTTPException, Query
from pydantic import BaseModel
from typing import List
import os
import requests
from pptx import Presentation
from PIL import Image
import base64
import re

# 创建 FastAPI app
app = FastAPI()

# 配置路径
UPLOAD_DIR = "/home/guwei/pptjiexi/ppt"
IMAGE_DIR = "/home/guwei/pptjiexi/image"
TXT_DIR = "/home/guwei/pptjiexi/output"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(IMAGE_DIR, exist_ok=True)
os.makedirs(TXT_DIR, exist_ok=True)

# DeepSeek-VL配置
DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"
DEEPSEEK_API_KEY = "sk-63ba35c4ee704383b23249a89617bdc2"

# 定义返回结构
class PageSpeech(BaseModel):
    page_index: int
    speech_script: str

class VisualUnderstandResponse(BaseModel):
    speeches: List[PageSpeech]

def extract_text_from_slide(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            texts.append(shape.text)
    return "\n".join(texts)

def extract_images_from_slide(slide, idx) -> str:
    img = Image.new("RGB", (1280, 720), color="white")
    img_path = os.path.join(IMAGE_DIR, f"slide_{idx + 1}.jpg")
    img.save(img_path)
    return img_path

def generate_prompt(ppt_text: str, language: str) -> str:
    """
    根据目标语言生成不同的提示词
    """
    if language.lower() == "en":
        # 英文提示
        prompt = f"""
You are a professional lecturer. Please generate a natural, smooth, content-rich speech based on the slide content.
Requirements:
- Avoid template openings like "Okay, next we will see..." or "Let's look at...".
- Dive directly into the topic in a natural, spoken style.
- Explain the texts, images, and formulas in depth.
- Control each slide's speech to around 30-60 seconds.
- Ensure smooth transitions and clear logical structure.

Here is the slide content:

Slide text:
{ppt_text}

Slide image description:
Please infer based on the image content.

Please generate the speech text:
"""
    else:
        # 默认中文提示
        prompt = f"""
你是一名专业讲师，需要根据幻灯片内容，生成一段正式自然、内容丰富的演讲稿。
要求：
- 不要出现"好的，接下来"、"让我们来看一下"等模板化开场。
- 直接进入主题讲解，保持口语化。
- 结合文字、图片、公式深入讲解。
- 每页时长控制30秒到60秒，保证信息密度。
- 用流畅自然的过渡衔接上下文。
- 文字清晰有条理，避免空洞描述。

下面是幻灯片内容：

幻灯片文字内容：
{ppt_text}

幻灯片图片描述：
请根据图片推理补充讲解。

请直接输出正式演讲稿内容：
"""
    return prompt

def call_deepseek_chat(ppt_text: str, image_base64: str, language: str) -> str:
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }

    prompt = generate_prompt(ppt_text, language)

    body = {
        "model": "deepseek-chat",
        "messages": [
            {"role": "system", "content": "You are a professional assistant for generating lecture scripts."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7
    }

    response = requests.post(DEEPSEEK_API_URL, headers=headers, json=body)

    if response.status_code != 200:
        raise Exception(f"DeepSeek调用失败: {response.text}")

    data = response.json()
    speech_text = data['choices'][0]['message']['content']
    return speech_text

def clean_speech_text(text: str) -> str:
    """
    清理括号提示和常见开场
    """
    text = re.sub(r"（.*?）", "", text)
    text = re.sub(r"\(.*?\)", "", text)
    text = re.sub(r"^(好的|那么|接下来|下面)[，,。]*", "", text.strip())
    return text.strip()

@app.post("/api/visual-understanding", response_model=VisualUnderstandResponse)
async def visual_understanding(
    file: UploadFile = File(...),
    language: str = Query(default="zh", description="输出语言，如 zh 或 en")
):
    """
    生成每页演讲稿，支持多语言输出
    """
    if not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="只支持上传PPTX文件。")

    pptx_path = os.path.join(UPLOAD_DIR, file.filename)
    with open(pptx_path, "wb") as f:
        f.write(await file.read())

    prs = Presentation(pptx_path)

    speeches = []
    txt_content = ""

    for idx, slide in enumerate(prs.slides):
        ppt_text = extract_text_from_slide(slide)
        img_path = extract_images_from_slide(slide, idx)

        with open(img_path, "rb") as f:
            image_base64 = base64.b64encode(f.read()).decode('utf-8')

        try:
            raw_speech_text = call_deepseek_chat(ppt_text, image_base64, language)
            speech_text = clean_speech_text(raw_speech_text)
        except Exception as e:
            speech_text = f"第{idx + 1}页生成失败: {str(e)}"

        speeches.append(PageSpeech(page_index=idx + 1, speech_script=speech_text))

        # 写入TXT
        txt_content += f"\n=== 第{idx + 1}页 ===\n{speech_text}\n"

    # 保存最终txt
    txt_filename = os.path.splitext(file.filename)[0] + f"_speech_{language}.txt"
    txt_path = os.path.join(TXT_DIR, txt_filename)
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(txt_content)

    return VisualUnderstandResponse(speeches=speeches)
